from __future__ import annotations
from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pydantic import BaseModel
from utils.fs import SLIDES_OUT, ts

# -----------------------
# Models
# -----------------------
class Bullet(BaseModel):
    text: str
    timestamp: str | None = None

class Section(BaseModel):
    heading: str
    bullets: list[Bullet]

class Outline(BaseModel):
    title: str
    topics: list[str]
    sections: list[Section]

# Slide layouts (standard 4:3)
TITLE_LAYOUT = 0
TITLE_AND_CONTENT = 1
SECTION_HEADER = 2

# -----------------------
# Color Scheme (Clean Academic Blue)
# -----------------------
COLORS = {
    "primary_blue": RGBColor(59, 130, 246),      # #3b82f6
    "dark_blue": RGBColor(30, 64, 175),          # #1e40af
    "light_blue": RGBColor(239, 246, 255),       # #eff6ff
    "white": RGBColor(255, 255, 255),
    "dark_gray": RGBColor(55, 65, 81),           # #374151
    "medium_gray": RGBColor(107, 114, 128),      # #6b7280
}

# -----------------------
# Font Settings
# -----------------------
FONTS = {
    "title_size": Pt(44),
    "subtitle_size": Pt(24),
    "heading_size": Pt(32),
    "bullet_size": Pt(20),
    "small_text": Pt(14),
}

# -----------------------
# Helper Functions
# -----------------------
def set_clean_background(slide):
    """Set a clean white background with light blue accent"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["white"]

def format_title_safe(shape, text: str, max_length: int = 80):
    """Safely format title text with length limit"""
    shape.text = text[:max_length] + ("..." if len(text) > max_length else "")
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER

def format_heading_safe(shape, text: str, max_length: int = 60):
    """Safely format heading text with length limit"""
    shape.text = text[:max_length] + ("..." if len(text) > max_length else "")
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.color.rgb = COLORS["dark_blue"]

def create_bullet_safe(text_frame, text: str, max_length: int = 120):
    """Create a bullet point with safe length handling"""
    clean_text = text[:max_length] + ("..." if len(text) > max_length else "")
    if text_frame.text == "":
        paragraph = text_frame.paragraphs[0]
    else:
        paragraph = text_frame.add_paragraph()
    
    paragraph.text = clean_text
    paragraph.level = 0
    paragraph.font.size = FONTS["bullet_size"]
    paragraph.font.color.rgb = COLORS["dark_gray"]
    paragraph.space_after = Pt(8)

# -----------------------
# Slide Creation Functions
# -----------------------
def _add_title_slide(prs: Presentation, title: str, topics: list[str]):
    """Create a clean, professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    set_clean_background(slide)
    
    # Format title
    title_shape = slide.shapes.title
    format_title_safe(title_shape, title)
    title_shape.text_frame.paragraphs[0].font.size = FONTS["title_size"]
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS["dark_blue"]
    
    # Format subtitle with topics
    if slide.placeholders and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        topics_text = " • ".join(topics[:3])  # Limit to 3 main topics
        subtitle_shape.text = topics_text[:100]  # Safe length
        for paragraph in subtitle_shape.text_frame.paragraphs:
            paragraph.font.size = FONTS["subtitle_size"]
            paragraph.font.color.rgb = COLORS["medium_gray"]
            paragraph.font.italic = True
            paragraph.alignment = PP_ALIGN.CENTER

def _add_section_slide(prs: Presentation, section: Section):
    """Create a clean content slide for a section"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    set_clean_background(slide)
    
    # Format heading
    title_shape = slide.shapes.title
    format_heading_safe(title_shape, section.heading)
    title_shape.text_frame.paragraphs[0].font.size = FONTS["heading_size"]
    
    # Format bullet points
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()  # Clear default text
    
    # Add bullet points (limit to 6 per slide for readability)
    for bullet in section.bullets[:6]:
        create_bullet_safe(text_frame, bullet.text)
    
    # Add continuation note if bullets were truncated
    if len(section.bullets) > 6:
        continuation_para = text_frame.add_paragraph()
        continuation_para.text = f"... and {len(section.bullets) - 6} more points"
        continuation_para.font.size = FONTS["small_text"]
        continuation_para.font.color.rgb = COLORS["medium_gray"]
        continuation_para.font.italic = True

def _add_topics_overview_slide(prs: Presentation, topics: list[str]):
    """Create a clean overview slide for multiple topics"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    set_clean_background(slide)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = "Lecture Topics"
    title_shape.text_frame.paragraphs[0].font.size = FONTS["heading_size"]
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS["dark_blue"]
    
    # Topics as bullet points
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for topic in topics[:8]:  # Limit topics for readability
        create_bullet_safe(text_frame, topic)
    
    if len(topics) > 8:
        continuation_para = text_frame.add_paragraph()
        continuation_para.text = f"... and {len(topics) - 8} more topics"
        continuation_para.font.size = FONTS["small_text"]
        continuation_para.font.color.rgb = COLORS["medium_gray"]
        continuation_para.font.italic = True

def _add_agenda_slide(prs: Presentation, sections: list[Section]):
    """Create a clean agenda slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    set_clean_background(slide)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = "Lecture Agenda"
    title_shape.text_frame.paragraphs[0].font.size = FONTS["heading_size"]
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS["dark_blue"]
    
    # Section headings as agenda items
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for i, section in enumerate(sections[:10], 1):  # Limit to 10 sections
        agenda_item = f"{i}. {section.heading}"
        create_bullet_safe(text_frame, agenda_item)

# -----------------------
# Main Function
# -----------------------
def outline_to_pptx(outline: Outline, filename_stem: Optional[str] = None) -> Path:
    """Convert outline to clean, professional PowerPoint presentation"""
    prs = Presentation()
    
    # Standard 4:3 slide size (default)
    
    # 1. Title Slide
    _add_title_slide(prs, outline.title, outline.topics)
    
    # 2. Agenda Slide (if multiple sections)
    if len(outline.sections) > 1:
        _add_agenda_slide(prs, outline.sections)
    
    # 3. Topics Overview (if multiple topics)
    if len(outline.topics) > 1:
        _add_topics_overview_slide(prs, outline.topics)
    
    # 4. Content Slides
    for section in outline.sections:
        _add_section_slide(prs, section)
    
    # 5. Summary Slide
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
    set_clean_background(slide)
    
    title_shape = slide.shapes.title
    title_shape.text = "Thank You"
    title_shape.text_frame.paragraphs[0].font.size = FONTS["title_size"]
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS["dark_blue"]
    
    if slide.placeholders and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = "Questions?"
        subtitle_shape.text_frame.paragraphs[0].font.size = FONTS["subtitle_size"]
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = COLORS["medium_gray"]
    
    # Save presentation
    stem = filename_stem or f"lecture_slides_{ts()}"
    out_path = SLIDES_OUT / f"{stem}.pptx"
    prs.save(out_path)
    
    return out_path